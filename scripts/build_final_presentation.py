"""
scripts/build_final_presentation.py

reference/(최종) AI_Master_Project_최종발표_멘티명_사번.pptx 템플릿의 서식(폰트/색상/굵기/
레이아웃)을 그대로 유지한 채, 지금까지 작성된 코드와 산출물 문서로 채울 수 있는
placeholder 텍스트만 실제 내용으로 교체한다.

원칙:
- Bold로 된 문장은 대부분 템플릿이 고정으로 두는 "질문/라벨"이라 손대지 않는다
  (예: "선택 이유:", "핵심 활용:", "어떻게 해결했는가? (설계 결정)").
- Non-bold "예: ..." 식 설명/예시 텍스트만 실제 내용으로 교체한다.
- 멘티 성명/사번, 멘토 성명처럼 코드·산출물 어디에도 없는 정보는 절대 채우지 않고
  원본 템플릿 placeholder([성명], [사번] 등) 그대로 둔다.
- 정량적 수치는 이번 세션에서 실제로 측정한 것(Meta Search 병렬화 처리시간 단축)만 쓰고,
  실측하지 않은 KPI(목표 수치일 뿐인 것)는 채우지 않는다.
"""

import os
from pptx import Presentation
from pptx.util import Inches

SRC = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "reference", "(최종) AI_Master_Project_최종발표_멘티명_사번.pptx",
)
OUT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "산출물", "AI_Master_Project_최종발표_0826_v2.pptx",
)


def set_run(shape, p_idx, r_idx, text):
    shape.text_frame.paragraphs[p_idx].runs[r_idx].text = text


def clear_runs(shape, p_idx, r_indices):
    for r_idx in r_indices:
        set_run(shape, p_idx, r_idx, "")


def main():
    prs = Presentation(SRC)
    s1, s2, s3, s4 = prs.slides[0], prs.slides[1], prs.slides[2], prs.slides[3]

    # ── 슬라이드 1: 표지 ──────────────────────────────────────
    # 과제명만 채움. 멘티/멘토 정보는 알 수 없어 그대로 둔다.
    set_run(s1.shapes[4], 0, 1, "데이터 명세관(SchemaScout) — 명세서·실데이터 정합성 자동 판별 AI 에이전트")

    # ── 슬라이드 2: 프로젝트 개요 ─────────────────────────────
    # 1. 문제 정의
    box = s2.shapes[10]
    set_run(box, 1, 0, "데이터 판매 사업부 담당자가 수십~수천 개 컬럼·수십 개월 요구기간의 명세서를 "
                       "실 DB와 일일이 대조해 존재·타입·기간을 전수 수작업으로 확인하던 업무를 자동화")
    set_run(box, 4, 0, "판단 기준이 문서화되지 않아 특정 담당자에게 의존하고, DataLake 미적재 구간을 놓치면 "
                       "제공 불가 데이터를 제공 가능으로 잘못 안내해 신뢰 리스크가 발생하며, 검증 지연은 "
                       "고객 회신 지연·영업 기회 손실로 이어짐")

    # 2. 핵심 기능
    box = s2.shapes[14]
    set_run(box, 1, 0, "RAG 기반 컬럼 매칭 + Human-in-the-Loop를 결합한 6-Agent LangGraph 파이프라인")
    set_run(box, 4, 0, "Parsing → Meta Search → Join Resolution → DB Validation → Classification → Report")
    set_run(box, 8, 0, "LangGraph / DuckDB vss(RAG) / Azure OpenAI Structured Output / FastAPI+Streamlit")

    # 핵심 성과 - 수치1(정성적): 채움 / 수치2(정량적): 실측치로 채움 / 수치3: 실측 없어 템플릿 유지
    box = s2.shapes[17]
    set_run(box, 0, 0, "판별 로직 표준화")
    clear_runs(box, 1, [1, 2, 3])
    set_run(box, 1, 0, "암묵지에 의존하던 수작업 판별을 재현·설명 가능한 규칙+LLM 로직으로 표준화(6-Agent + HITL 7종)")

    box = s2.shapes[19]
    set_run(box, 0, 0, "약 17%")
    clear_runs(box, 1, [1, 2, 3])
    set_run(box, 1, 0, "Meta Search Agent 병렬 처리 적용 시 근접매칭 컬럼 처리시간 단축(실측, 5컬럼 기준 31.02초→25.90초)")
    # 수치 3(slide2.shapes[21])은 실측된 업무 자동화율/검색시간 단축 지표가 없어 템플릿 그대로 둠

    # Key Message
    box = s2.shapes[24]
    set_run(box, 0, 0, "\"규칙 기반 매칭·LLM 판단·담당자 확인(HITL)을 결합해 컬럼 검증을 자동화하고,")
    set_run(box, 1, 0, "여러 테이블 간 조인 가능성까지 검증하는 6-Agent 파이프라인을 구축했습니다.\"")

    # ── 슬라이드 3: 기술 아키텍처 ─────────────────────────────
    # 좌측 카드(Shape 7: L=0.28" T=0.73" W=5.05" H=4.67")에 실제 다이어그램 이미지를 삽입하고,
    # 텍스트 placeholder(제목/예시 불릿)는 이미지가 그 내용을 대신하므로 전부 비운다.
    box = s3.shapes[8]
    for p_idx in range(len(box.text_frame.paragraphs)):
        for r_idx in range(len(box.text_frame.paragraphs[p_idx].runs)):
            set_run(box, p_idx, r_idx, "")

    diagram_path = os.path.join(os.path.dirname(OUT), "architecture_diagram.png")
    if os.path.exists(diagram_path):
        from PIL import Image
        with Image.open(diagram_path) as im:
            img_w_px, img_h_px = im.size
        card_left_in, card_top_in, card_w_in, card_h_in = 0.28, 0.73, 5.05, 4.67
        pad_in = 0.15
        max_w_in = card_w_in - 2 * pad_in
        max_h_in = card_h_in - 2 * pad_in
        aspect = img_w_px / img_h_px
        draw_w_in = max_w_in
        draw_h_in = draw_w_in / aspect
        if draw_h_in > max_h_in:
            draw_h_in = max_h_in
            draw_w_in = draw_h_in * aspect
        left_in = card_left_in + (card_w_in - draw_w_in) / 2
        top_in = card_top_in + (card_h_in - draw_h_in) / 2
        s3.shapes.add_picture(
            diagram_path,
            Inches(left_in), Inches(top_in), Inches(draw_w_in), Inches(draw_h_in),
        )

    # 기술 01: LangGraph + HITL
    set_run(s3.shapes[14], 0, 0, "LangGraph + Human-in-the-Loop(interrupt())")
    box = s3.shapes[15]
    set_run(box, 1, 0, "애매한 판단(유사도 매칭, 타입 불일치, 동일 컬럼명 중복, 조인키 추정)마다 사람이 개입할 "
                       "지점을 표현하려면 조건부 분기(add_conditional_edges)와 확인 후 재개(interrupt())가 "
                       "가능한 그래프 구조가 필요했음")
    set_run(box, 4, 0, "6-Agent + Meta Search 내부 7개 서브노드를 하나의 StateGraph로 구성, SqliteSaver "
                       "체크포인터로 확인 대기 상태를 세션 재시작 후에도 보존")

    # 기술 02: RAG + Structured Output
    set_run(s3.shapes[20], 0, 0, "RAG(DuckDB vss) + Structured Output(Pydantic)")
    box = s3.shapes[21]
    set_run(box, 1, 0, "컬럼명 표기가 다양해 규칙만으론 매칭이 안 되는 케이스가 많아 유사도 검색이 필요했고, "
                       "LLM 응답이 스키마를 벗어나지 않도록 JSON 모드 대신 Pydantic 기반 Structured Output으로 강제")
    set_run(box, 4, 0, "동일 컬럼명이 여러 테이블에 존재하는 경우(예: 6개 테이블 공통 존재)를 감지해 임의 선택 "
                       "대신 담당자 확인으로 전환, 조인키 후보도 임베딩만으론 안 되고 실측 값 overlap(SEMI JOIN)"
                       "으로 이중 검증")

    # 기술 03: Azure OpenAI 모델 분리 + 재시도
    set_run(s3.shapes[26], 0, 0, "Azure OpenAI(gpt-4.1-mini / gpt-5-mini) — 모델 분리 + 재시도")
    box = s3.shapes[27]
    set_run(box, 1, 0, "빈도 높은 단순 분류(헤더 매핑)는 경량 모델로 비용을 아끼고, 의미 기반 추론이 핵심인 "
                       "매칭 판단에만 상대적으로 성능 좋은 모델을 배정. 실 운영 환경의 일시적 API 오류에도 "
                       "대비가 필요했음")
    set_run(box, 4, 0, "tenacity 지수 백오프(RateLimitError/APITimeoutError 등만 재시도, 영구 오류는 즉시 전파), "
                       "담당자 확인이 불필요한 구간만 선별해 ThreadPoolExecutor 병렬화")

    # ── 슬라이드 4: 핵심 기술 과제 ────────────────────────────
    box = s4.shapes[9]
    set_run(box, 0, 0, "이름·임베딩 유사도만으로 조인키를 확정하면, 의미상 유사해도 실제로는 값 컬럼이거나 "
                       "카디널리티가 안 맞는 컬럼까지 조인키로 오인될 위험이 있었습니다.")
    set_run(box, 1, 0, "예: total_ic_mou/total_og_mou(임베딩 유사도 0.84)는 값 컬럼이라 조인키가 아니었고, "
                       "aon(가입기간)은 값이 100% 겹쳐도 relation_type이 N:N이라 실제로는 저품질 키였음을 "
                       "실측으로 확인했습니다.")

    box = s4.shapes[13]
    set_run(box, 1, 0, "이름 일치 → 임베딩 유사도로 후보를 넓힌 뒤, 실제 데이터에 SEMI JOIN을 돌려 양방향 값 "
                       "포함률(containment)과 relation_type(1:1/1:N/N:N)을 실측한 후보만 채택. 확신도가 아무리 "
                       "높아도 새로 추정한 조인키는 항상 담당자 확인(interrupt())을 거치게 설계")
    set_run(box, 4, 0, "그래프 BFS(2-hop 조인 경로 탐색) + SEMI JOIN 기반 inclusion-dependency 검증 — 그라운드 "
                       "트루스(실제 값 overlap)가 존재하는 문제라 통계적 검증이 LLM 판단보다 신뢰할 수 있다고 "
                       "판단해 고급 추론 기법 대신 채택")
    set_run(box, 7, 0, "이름·설명 유사도만으로는 aon(가입기간, 실제로는 55,237개 중복값)과 mobile_number(진짜 "
                       "식별자)가 confidence 동점으로 나와 구분이 안 됐음을 실측으로 확인했기 때문")
    # "결과 및 성과(수치로 제시)" 3개 슬롯은 이 hurdle에 대한 정량적 개선치(오탐률 등)를 실측하지 않아 템플릿 유지

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"작성 완료: {OUT}")


if __name__ == "__main__":
    main()
